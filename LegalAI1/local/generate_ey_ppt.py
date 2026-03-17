from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# EY Brand Colors
EY_YELLOW = RGBColor(255, 230, 0)
EY_DARK_GREY = RGBColor(46, 46, 56)
WHITE = RGBColor(255, 255, 255)

def apply_ey_style(slide, title_text, body_text_list):
    # Set background color to EY Dark Grey for contrast, or keep it white with Dark Grey text
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = EY_DARK_GREY
    
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.color.rgb = EY_YELLOW
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    
    if hasattr(slide.shapes, 'placeholders') and len(slide.shapes.placeholders) > 1:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for text in body_text_list:
            p = tf.add_paragraph()
            p.text = text
            p.font.color.rgb = WHITE
            p.font.name = 'Arial'
            p.font.size = Pt(20)

def create_presentation():
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = EY_DARK_GREY
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Legal AI Platform"
    title.text_frame.paragraphs[0].font.color.rgb = EY_YELLOW
    title.text_frame.paragraphs[0].font.name = 'Arial'
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Next-Generation Document Intelligence & Contract Analysis\n\nBuilding a better working world"
    subtitle.text_frame.paragraphs[0].font.color.rgb = WHITE
    subtitle.text_frame.paragraphs[0].font.name = 'Arial'

    # Content Slides Data
    slides_data = [
        ("Executive Summary", 
         ["• The Challenge: Legal teams spend thousands of hours manually reviewing contracts.",
          "• The Solution: Automated platform that ingests, comprehends, and queries complex legal documents.",
          "• The Value: Reduce review time by 80%, mitigate risk, and scale seamlessly."]),
        ("Core Capabilities",
         ["• Automated Data Extraction: LlamaParse integration for precise markdown conversion.",
          "• Intelligent Summarization: AI instantly generates concise executive summaries.",
          "• Risk Analysis: Categorizes legal and compliance risks directly from the text.",
          "• Matter Management: Organizes documents securely into hierarchical structures."]),
        ("Advanced Search & Q&A",
         ["• Hybrid Search Engine: Combines keyword matching with semantic pgvector search.",
          "• Q&A RAG: Context-aware Chat functionality for natural language legal questions.",
          "• Direct Citations: Traces AI responses back to exact source paragraphs."]),
        ("Enterprise Architecture",
         ["• Frontend/Backend: Next.js & React",
          "• Database: PostgreSQL with pgvector & Drizzle ORM",
          "• Asynchronous Pipeline: Redis & background Python workers",
          "• AI Models: Multi-LLM integration for varied legal tasks"]),
        ("Security & Compliance",
         ["• Strict Tenant Isolation: Segregated by Organization ID",
          "• Audit Logging: Transparent tracking of document processing",
          "• RBAC: Granular permissions mapped via OAuth"]),
        ("Roadmap",
         ["• Adverse Media Screening API Integration",
          "• Multi-lingual Contract Analysis",
          "• Advanced Graph Relationships for Statutes"])
    ]
    
    # Add slides dynamically
    bullet_slide_layout = prs.slide_layouts[1]
    for title_text, body_text in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        apply_ey_style(slide, title_text, body_text)
        
    prs.save('Legal_AI_EY_Presentation.pptx')
    print("Presentation saved as 'Legal_AI_EY_Presentation.pptx'")

if __name__ == '__main__':
    create_presentation()
